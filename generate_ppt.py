#!/usr/bin/env python3
"""Generate Furniture Sale PowerPoint from selling items."""

import io
import sys
import requests
from bs4 import BeautifulSoup
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE

# ─────────────────────────────────────────────────────────────
# ITEMS  (name, online price, product link or None)
# ─────────────────────────────────────────────────────────────
ITEMS = [
    {"name": "Bed frame & mattress",      "price": "$500",   "link": None},
    {"name": "Mirror",                     "price": None,     "link": None},
    {"name": "Silver standing lamp",       "price": None,     "link": None},
    {"name": "Rattan hamper",              "price": None,     "link": None},
    {"name": "Black storage box",          "price": None,     "link": None},
    {"name": "Clothes organizer",          "price": None,     "link": None},
    {"name": "White & gold side table",    "price": "$32",    "link": None},
    {"name": "White desk lamp",            "price": "$21.35", "link": "https://www.amazon.com/lianheng-Outlet-Office-College-Flexible/dp/B0CQ81WRKD/"},
    {"name": "Honeywell fan",              "price": "$14.99", "link": "https://www.target.com/p/honeywell-turbo-force-table-air-circulator-fan-white/-/A-11153866"},
    {"name": "TV stand / console",         "price": "$100",   "link": None},
    {"name": "Insignia TV",                "price": "$79.99", "link": None},
    {"name": "Sun lamp",                   "price": "$79.99", "link": "https://www.amazon.com/dp/B075XRZTV1"},
    {"name": "Fluffy pillows",             "price": "$40",    "link": "https://www.target.com/p/euro-faux-mongolian-fur-decorative-throw-pillow-cream-threshold-8482/-/A-85376335"},
    {"name": "Desk and chair",             "price": "$115",   "link": None},
    {"name": "Weighing scale",             "price": "$13.99", "link": "https://www.amazon.com/dp/B08XXNDVFB"},
    {"name": "Drawer",                     "price": "$20",    "link": "https://www.target.com/p/3-drawer-wide-cart-white-brightroom-8482/-/A-84242464"},
    {"name": "2 bulletin boards",          "price": "$29.99", "link": None},
    {"name": "Toothbrush organizer",       "price": "$13",    "link": "https://www.target.com/p/accent-vanity-organizer-silver-threshold-8482/-/A-78367670"},
    {"name": "Bathroom mats",              "price": "$20",    "link": "https://www.target.com/p/2pk-quick-dry-bath-rug-set-threshold/-/A-81909483"},
    {"name": "Trashcan",                   "price": "$12",    "link": "https://www.target.com/p/2-7gal-step-trash-can-clear-brightroom-8482/-/A-87299507"},
    {"name": "Under bed shoe organizers",  "price": "$49.99", "link": "https://www.amazon.com/gp/aw/d/B0B6G1X86L/"},
    {"name": "Shampoo/conditioner holder", "price": "$10",    "link": "https://www.amazon.com/SUNFICON-Adhesive-Organizer-Essentials-Traceless/dp/B07HP1G8V6/"},
    {"name": "Floor pillow",              "price": "$59.99", "link": "https://www.amazon.com/dp/B0B195QT96"},
]

# ─────────────────────────────────────────────────────────────
# IMAGE HELPERS
# ─────────────────────────────────────────────────────────────
HEADERS = {
    "User-Agent": (
        "Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7) "
        "AppleWebKit/537.36 (KHTML, like Gecko) "
        "Chrome/124.0.0.0 Safari/537.36"
    ),
    "Accept": "text/html,application/xhtml+xml,application/xml;q=0.9,image/avif,image/webp,*/*;q=0.8",
    "Accept-Language": "en-US,en;q=0.9",
    "Accept-Encoding": "gzip, deflate, br",
}


def fetch_og_image_url(page_url):
    """Return the og:image URL from a product page, or None on failure."""
    try:
        r = requests.get(page_url, headers=HEADERS, timeout=20, allow_redirects=True)
        soup = BeautifulSoup(r.text, "lxml")

        # 1. og:image meta tag (most reliable)
        for tag in soup.find_all("meta", property="og:image"):
            src = tag.get("content", "").strip()
            if src.startswith("http"):
                return src

        # 2. Amazon landing image
        for tag in soup.find_all("img", id="landingImage"):
            src = tag.get("src") or tag.get("data-old-hires") or tag.get("data-src")
            if src and src.startswith("http"):
                return src

        # 3. Target product image JSON-LD
        for script in soup.find_all("script", type="application/ld+json"):
            text = script.string or ""
            if '"image"' in text:
                import json, re
                try:
                    data = json.loads(text)
                    imgs = data.get("image") if isinstance(data, dict) else None
                    if isinstance(imgs, list) and imgs:
                        return imgs[0]
                    if isinstance(imgs, str) and imgs.startswith("http"):
                        return imgs
                except Exception:
                    pass

    except Exception as e:
        print(f"    [warn] page fetch failed for {page_url}: {e}")
    return None


def download_image_buf(img_url):
    """Download an image URL and return a square-cropped PNG BytesIO, or None."""
    try:
        r = requests.get(img_url, headers=HEADERS, timeout=20)
        if r.status_code == 200 and r.content:
            img = Image.open(io.BytesIO(r.content)).convert("RGBA")
            # Paste onto white background (handles transparent PNGs)
            bg = Image.new("RGB", img.size, (255, 255, 255))
            if img.mode == "RGBA":
                bg.paste(img, mask=img.split()[3])
            else:
                bg = img.convert("RGB")
            # Square-crop from centre
            w, h = bg.size
            side = min(w, h)
            left = (w - side) // 2
            top = (h - side) // 2
            bg = bg.crop((left, top, left + side, top + side))
            buf = io.BytesIO()
            bg.save(buf, format="PNG")
            buf.seek(0)
            return buf
    except Exception as e:
        print(f"    [warn] image download failed: {e}")
    return None


def placeholder_buf(size=400):
    """Light gray placeholder square."""
    img = Image.new("RGB", (size, size), (235, 235, 235))
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    buf.seek(0)
    return buf


# ─────────────────────────────────────────────────────────────
# SLIDE BUILDERS
# ─────────────────────────────────────────────────────────────
SLIDE_W = Inches(10)
SLIDE_H = Inches(7.5)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
BLACK   = RGBColor(0x1A, 0x1A, 0x1A)
GRAY    = RGBColor(0x88, 0x88, 0x88)
LINE_C  = RGBColor(0xCC, 0xCC, 0xCC)
FONT    = "Calibri"


def set_white_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE


def add_text(slide, text, x, y, w, h, size, bold=False, color=BLACK, align=PP_ALIGN.CENTER, wrap=True):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = FONT


def build_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_white_bg(slide)

    # Main title
    add_text(slide, "Furniture Sale",
             Inches(1), Inches(2.6), Inches(8), Inches(1.8),
             size=52, bold=False, color=BLACK)

    # Sub-line
    add_text(slide, "All items priced to sell  •  DM to claim",
             Inches(1), Inches(4.5), Inches(8), Inches(0.7),
             size=15, color=GRAY)

    # Thin bottom rule
    line = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(1), Inches(5.4), Inches(8), Pt(1.5)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = LINE_C
    line.line.fill.background()


def build_items_slide(prs, row):
    """row: list of 1-3 item dicts (with '_img_buf' key added)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_white_bg(slide)

    n = len(row)
    margin   = Inches(0.5)
    gap      = Inches(0.35)
    col_w    = (SLIDE_W - 2 * margin - (n - 1) * gap) / n
    img_h    = Inches(4.0)
    top_img  = Inches(0.8)
    top_name = top_img + img_h + Inches(0.22)
    top_price= top_name + Inches(0.58)

    for i, item in enumerate(row):
        x = margin + i * (col_w + gap)

        # ── image / placeholder ──
        has_image = bool(item.get("_img_buf"))
        buf = item.get("_img_buf") or placeholder_buf()
        try:
            slide.shapes.add_picture(buf, x, top_img, col_w, img_h)
        except Exception:
            has_image = False

        if not has_image:
            # Draw border box over the placeholder
            rect = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, top_img, col_w, img_h)
            rect.fill.solid()
            rect.fill.fore_color.rgb = RGBColor(0xEB, 0xEB, 0xEB)
            rect.line.color.rgb = RGBColor(0xBB, 0xBB, 0xBB)
            # "Add photo" label inside placeholder
            add_text(slide, "[ Add photo ]",
                     x, top_img + img_h / 2 - Inches(0.25), col_w, Inches(0.5),
                     size=11, color=RGBColor(0xAA, 0xAA, 0xAA))

        # ── separator line ──
        sep = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            x, top_img + img_h + Inches(0.1), col_w, Pt(1)
        )
        sep.fill.solid()
        sep.fill.fore_color.rgb = LINE_C
        sep.line.fill.background()

        # ── item name ──
        add_text(slide, item["name"],
                 x, top_name, col_w, Inches(0.65),
                 size=13, color=BLACK)

        # ── price ──
        price = item["price"] if item["price"] else "Price on request"
        add_text(slide, price,
                 x, top_price, col_w, Inches(0.5),
                 size=13, color=RGBColor(0x44, 0x44, 0x44))


# ─────────────────────────────────────────────────────────────
# MAIN
# ─────────────────────────────────────────────────────────────
def build_notes_slide(prs, items_needing_photos):
    """Final slide listing all items that still need a photo added."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_white_bg(slide)

    add_text(slide, "Photos still needed",
             Inches(0.5), Inches(0.3), Inches(9), Inches(0.7),
             size=20, bold=True, color=BLACK, align=PP_ALIGN.LEFT)

    sep = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0.5), Inches(1.05), Inches(9), Pt(1)
    )
    sep.fill.solid()
    sep.fill.fore_color.rgb = LINE_C
    sep.line.fill.background()

    lines = []
    for item in items_needing_photos:
        link_note = f"  →  {item['link']}" if item.get("link") else "  →  (no link – add your own photo)"
        lines.append(f"• {item['name']}{link_note}")

    body = "\n".join(lines)
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(6))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = body
    run.font.size = Pt(9)
    run.font.color.rgb = RGBColor(0x44, 0x44, 0x44)
    run.font.name = FONT


def main():
    output = "/home/user/game1/furniture_sale.pptx"

    # 1. Fetch images (network may block; placeholders used as fallback)
    print(f"\nFetching images for {len(ITEMS)} items...\n")
    for item in ITEMS:
        link = item.get("link")
        if link:
            print(f"  [{item['name']}]")
            og_url = fetch_og_image_url(link)
            if og_url:
                print(f"    downloading image...")
                item["_img_buf"] = download_image_buf(og_url)
            else:
                print(f"    no image found – placeholder")
                item["_img_buf"] = None
        else:
            item["_img_buf"] = None

    # 2. Build presentation
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    build_title_slide(prs)

    for i in range(0, len(ITEMS), 3):
        build_items_slide(prs, ITEMS[i:i + 3])

    # Notes slide for items missing photos
    no_photo = [it for it in ITEMS if not it.get("_img_buf")]
    if no_photo:
        build_notes_slide(prs, no_photo)

    prs.save(output)
    print(f"\n✓ Saved: {output}")
    total_slides = 1 + (len(ITEMS) + 2) // 3 + (1 if no_photo else 0)
    print(f"  {total_slides} slides total  ({len(no_photo)} items need photos added manually)")


if __name__ == "__main__":
    main()
